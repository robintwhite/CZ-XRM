from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import json
import argparse
import os
import pandas as pd

"""

Slide Master layout
[0] = TitleSlide
[1] = ContentSlide 1 image
[2] = JobDescSlide 4 images
[3] = SummarySlide 4 images
[4] = QuadViewSlide 4 images
[5] = ComparisonSlide 2 images
[6] = TitleOnly
[7] = ScanParamsSlide
[8] = AppendixSlide
[9] = BlankSlide

See slides.json for more info or ZEISS_blank.pptx

Use include_slides.txt to define layout
Use image names to define which images go in where
Use txt/csv or json to define scan params

"""

colours = {'red': (0xFF, 0x00, 0x00),
           'green': (0x00, 0xFF, 0x00),
           'blue': (0x00, 0x00, 0xFF)}


def read_master(filename):

    slide_dic = {}
    prs = Presentation(filename)
    # title_slide_layout = prs.slide_layouts[1]
    # slide = prs.slides.add_slide(title_slide_layout)
    for slide in prs.slide_layouts:
        slide_index = prs.slide_layouts.index(slide)
        slide_dic[slide_index] = []
        for shape in slide.placeholders:
            slide_dic[slide_index].append({shape.placeholder_format.idx: shape.name})

            print('%d %d %s' % (prs.slide_layouts.index(slide), shape.placeholder_format.idx, shape.name))

    with open('slides.json', 'w') as outfile:
        json.dump(slide_dic, outfile, indent=2)


def create_slides(prs):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[12]

    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"

    prs.save('test.pptx')


def num_img_to_slide_master(num_imgs):

    # number of images in a slide, slide index from master
    #Use read_master to check jason for correct template number or in base ppt

    switcher = {
        0: 0,
        1: 1,
        2: 5,
        4: 4,
    }
    # switcher = {
    #     0: 0,
    #     1: 1,
    #     2: 5,
    #     4: 4,
    # }

    return switcher.get(num_imgs, "Invalid value")


def read_slide_list():

    text_file = open("include_slides.txt", "r")
    slide_list = text_file.read().split(',')
    text_file.close()
    slide_list = list(map(int, slide_list))
    return slide_list

def adjust_picture_to_fit(picture):

    pos_left, pos_top = picture.left, picture.top

    available_width = picture.width
    available_height = picture.height

    image_width, image_height = picture.image.size

    placeholder_aspect_ratio = float(available_width) / float(available_height)
    image_aspect_ratio = float(image_width) / float(image_height)

    picture.crop_top = 0
    picture.crop_left = 0
    picture.crop_bottom = 0
    picture.crop_right = 0

    # ---if the placeholder is "wider" in aspect, shrink the picture width while
    # ---maintaining the image aspect ratio
    if placeholder_aspect_ratio > image_aspect_ratio:
        picture.width = int(image_aspect_ratio * available_height)
        picture.height = available_height
    # ---otherwise shrink the height
    else:
        picture.height = int(available_width / image_aspect_ratio)
        picture.width = available_width

    picture.left, picture.top = pos_left, pos_top
    #print(picture.left, picture.top)

def insert_images(slide, slide_num, images_path, image_df):

    """
    Insert images into a slide.
    :param slide: = slide object from Presentation class
    :param slide_num: the template slide number for formatting
    :param images_path: the directory to the folder with all the images
    :param image_df: Pandas data frame regarding information of each image in images_path
    :return: None
    """

    placeholders = get_image_placeholders(slide)
    print('placeholders: ', placeholders)
    image_pool = image_df[image_df['slide_num'] == slide_num]
    print('image_pool: ', image_pool)
    try:
        assert len(placeholders) == len(image_pool.index)
    except AssertionError:
        print('Length of placeholders in slide does not match image naming.')
    i = 0
    for idx, image in image_pool.iterrows():
        #print(image)
        image_path = os.path.join(images_path, image.path)
        placeholder = slide.placeholders[placeholders[i]]
        pic = placeholder.insert_picture(image_path)

        adjust_picture_to_fit(pic)

        #print(image.path)

        line = pic.line
        #print(image['view'])
        if image['view'] == 'red':
            line.color.rgb = RGBColor(255, 0, 0)
        elif image['view'] == 'green':
            line.color.rgb = RGBColor(0, 255, 0)
        elif image['view'] == 'blue':
            line.color.rgb = RGBColor(50, 205, 255)
        else:
            line.color.rgb = RGBColor(0, 0, 0)
        line.width = Pt(2.25)
        i+=1



def get_image_list(image_dir):

    image_names_list = []
    for img_filename in os.listdir(image_dir):
        if img_filename.endswith(".png") or img_filename.endswith(".PNG"):
            image_names_list.append(img_filename)
    return image_names_list


def create_image_df(image_names_list):

    image_list = []
    for img_filename in image_names_list:
        slide_num, view, name = img_filename.split('.')[0].split('_')
        image_list.append([int(slide_num), view, name])

    image_df = pd.DataFrame(image_list, columns=['slide_num', 'view', 'name'])
    image_df['path'] = pd.Series(image_names_list)
    #print(image_df)
    return image_df


def get_image_placeholders(slide):
    image_placeholder_list = []
    for placeholder in slide.placeholders:
        if "Picture" in placeholder.name:
            image_placeholder_list.append(placeholder.placeholder_format.idx)
    return image_placeholder_list


if __name__ == '__main__':

    # TODO: Get scan params using ZEISS python API for table

    cwdir = os.getcwd()
    parser = argparse.ArgumentParser()
    parser.add_argument('-d', '--images_dir',
                        required=True,
                        help="Directory containing the images.")
    parser.add_argument('-o', '--output_dir',
                        required=False,
                        default=os.getcwd(),
                        help="Output directory.")

    args = parser.parse_args()

    read_master(os.path.join(cwdir, 'Template\\ZEISS_blank.pptx'))

    images_path = args.images_dir #os.path.join(cwdir, 'Images')
    out_dir = args.output_dir
    image_names_list = get_image_list(images_path)

    prs = Presentation(os.path.join(cwdir, 'Template\\ZEISS_blank.pptx'))

    create_slides(prs)

    image_df = create_image_df(image_names_list)
    # Count number of slides required
    num_images_in_slide = image_df.groupby('slide_num')['path'].nunique()
    print(num_images_in_slide)
    for slide_num, num_imgs in num_images_in_slide.items():
        slide_master_num = num_img_to_slide_master(num_imgs)
        print('slidenum: {}, num_images: {}, slide_master_id: {}'.format(slide_num, num_imgs, slide_master_num))
        #print(slide_master_num)
        slide = prs.slides.add_slide(prs.slide_layouts[slide_master_num])
        insert_images(slide, slide_num, images_path, image_df)

    prs.save(os.path.join(out_dir, 'Draft.pptx'))
